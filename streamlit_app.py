# === VERSION FINAL TOTAL (ESTABLE + COMPLETA + EXPORTES) ===

import streamlit as st
import pandas as pd
import plotly.express as px
import matplotlib.pyplot as plt
import os

from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image
from reportlab.lib.styles import getSampleStyleSheet

from pptx import Presentation
from pptx.util import Inches

# -------------------------
# CONFIG
# -------------------------
st.set_page_config(page_title="Dashboard Ejecutivo", layout="wide")

if "vista" not in st.session_state:
    st.session_state.vista = "principal"

# -------------------------
# FUNCION GRAFICAS
# -------------------------
def guardar_img(nombre, plot_func):
    plt.figure(figsize=(10,5))
    plt.rcParams["figure.dpi"] = 120
    plot_func()
    plt.tight_layout()
    plt.savefig(nombre)
    plt.close()
    return nombre

# -------------------------
# PDF
# -------------------------
def generar_pdf(df_m, df, ventas, ganancia, margen, ratio):

    doc = SimpleDocTemplate("reporte_final.pdf")
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("REPORTE EJECUTIVO", styles['Title']))
    story.append(Spacer(1, 20))
    story.append(Paragraph("Análisis del negocio", styles['Normal']))
    story.append(PageBreak())

    story.append(Paragraph(f"Ventas: ${ventas:,.0f}", styles['Normal']))
    story.append(Paragraph(f"Ganancia: ${ganancia:,.0f}", styles['Normal']))
    story.append(Paragraph(f"Margen: {margen:.1f}%", styles['Normal']))
    story.append(PageBreak())

    img = guardar_img("temp_tend.png", lambda: (
        plt.plot(df_m["Periodo"], df_m["Ventas"], label="Ventas"),
        plt.plot(df_m["Periodo"], df_m["Ganancia"], label="Ganancia"),
        plt.legend(),
        plt.xticks(rotation=45)
    ))

    story.append(Paragraph("Tendencia", styles['Heading1']))
    story.append(Image(img, width=500, height=300))
    story.append(PageBreak())

    if ratio > 0.30:
        estado = "Alta volatilidad"
    elif ratio > 0.15:
        estado = "Media volatilidad"
    else:
        estado = "Estable"

    story.append(Paragraph(f"Estado: {estado}", styles['Normal']))

    doc.build(story)

    if os.path.exists("temp_tend.png"):
        os.remove("temp_tend.png")

# -------------------------
# PPT
# -------------------------
def generar_ppt(df_m, df, ventas, ganancia, margen):

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Reporte Ejecutivo"
    slide.placeholders[1].text = "Análisis del negocio"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "KPIs"
    slide.placeholders[1].text = f"""
Ventas: ${ventas:,.0f}
Ganancia: ${ganancia:,.0f}
Margen: {margen:.1f}%
"""

    img = guardar_img("ppt_temp.png", lambda: (
        plt.plot(df_m["Periodo"], df_m["Ventas"]),
        plt.plot(df_m["Periodo"], df_m["Ganancia"])
    ))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Tendencia"
    slide.shapes.add_picture(img, Inches(1), Inches(2), width=Inches(8))

    prs.save("reporte_final.pptx")

    if os.path.exists("ppt_temp.png"):
        os.remove("ppt_temp.png")

# -------------------------
# DATA
# -------------------------
archivo = st.file_uploader("Sube tu Excel", type=["xlsx"])

if archivo:

    df = pd.read_excel(archivo)
    df.columns = df.columns.str.strip()

    df["Fecha"] = pd.to_datetime(df["Fecha"], errors="coerce")
    df = df.dropna(subset=["Fecha"])

    df["Ganancia"] = df["Ventas"] - df["Costos"]
    df["Periodo"] = df["Fecha"].dt.to_period("M").astype(str)

    # -------------------------
    # FILTROS
    # -------------------------
    st.sidebar.header("Filtros")

    rango = st.sidebar.date_input("Fecha", [df["Fecha"].min(), df["Fecha"].max()])
    if len(rango) == 2:
        df = df[(df["Fecha"] >= pd.to_datetime(rango[0])) &
                (df["Fecha"] <= pd.to_datetime(rango[1]))]

    if "Pais" in df.columns:
        pais = st.sidebar.multiselect("País", df["Pais"].unique(), df["Pais"].unique())
        df = df[df["Pais"].isin(pais)]

    if "Region" in df.columns:
        reg = st.sidebar.multiselect("Región", df["Region"].unique(), df["Region"].unique())
        df = df[df["Region"].isin(reg)]

    if "Nombre" in df.columns:
        nom = st.sidebar.multiselect("Nombre", df["Nombre"].unique(), df["Nombre"].unique())
        df = df[df["Nombre"].isin(nom)]

    if df.empty:
        st.warning("Sin datos")
        st.stop()

    df_m = df.groupby("Periodo")[["Ventas", "Ganancia"]].sum().reset_index()

    ventas = df["Ventas"].sum()
    ganancia = df["Ganancia"].sum()
    margen = (ganancia / ventas * 100) if ventas != 0 else 0

    # =========================
    # PRINCIPAL
    # =========================
    if st.session_state.vista == "principal":

        st.title("📊 Dashboard Ejecutivo")

        c1, c2, c3 = st.columns(3)
        c1.metric("Ventas", f"${ventas:,.0f}")
        c2.metric("Ganancia", f"${ganancia:,.0f}")
        c3.metric("Margen", f"{margen:.1f}%")

        fig = px.line(df_m, x="Periodo", y=["Ventas", "Ganancia"], markers=True)
        st.plotly_chart(fig, use_container_width=True)

        col1, col2, col3, col4 = st.columns(4)

        if col1.button("🚦 Volatilidad"):
            st.session_state.vista = "volatilidad"

        if col2.button("👤 Responsables"):
            st.session_state.vista = "responsables"

        if col3.button("🔎 Causas"):
            st.session_state.vista = "causas"

        if col4.button("📄 Reporte"):
            st.session_state.vista = "reporte"

    # -------------------------
    # VOLATILIDAD
    # -------------------------
    elif st.session_state.vista == "volatilidad":

        if st.button("⬅️ Volver"):
            st.session_state.vista = "principal"

        media = df_m["Ventas"].mean()
        vol = df_m["Ventas"].std()
        ratio = vol / media if media != 0 else 0

        if ratio > 0.30:
            st.error("🔴 Alta volatilidad")
        elif ratio > 0.15:
            st.warning("🟡 Volatilidad media")
        else:
            st.success("🟢 Estable")

        st.line_chart(df_m.set_index("Periodo")["Ventas"])

    # -------------------------
    # RESPONSABLES
    # -------------------------
    elif st.session_state.vista == "responsables":

        if st.button("⬅️ Volver"):
            st.session_state.vista = "principal"

        df_nom = df.groupby(["Periodo", "Nombre"])["Ventas"].sum().reset_index()
        fig = px.line(df_nom, x="Periodo", y="Ventas", color="Nombre")
        st.plotly_chart(fig, use_container_width=True)

    # -------------------------
    # CAUSAS
    # -------------------------
    elif st.session_state.vista == "causas":

        if st.button("⬅️ Volver"):
            st.session_state.vista = "principal"

        df_reg = df.groupby(["Periodo", "Region"])["Ventas"].sum().reset_index()
        fig = px.line(df_reg, x="Periodo", y="Ventas", color="Region")
        st.plotly_chart(fig, use_container_width=True)

    # -------------------------
    # REPORTE
    # -------------------------
    elif st.session_state.vista == "reporte":

        if st.button("⬅️ Volver"):
            st.session_state.vista = "principal"

        media = df_m["Ventas"].mean()
        vol = df_m["Ventas"].std()
        ratio = vol / media if media != 0 else 0

        if st.button("📄 Generar PDF"):
            generar_pdf(df_m, df, ventas, ganancia, margen, ratio)
            with open("reporte_final.pdf", "rb") as f:
                st.download_button("Descargar PDF", f, "reporte_final.pdf")

        if st.button("📊 Generar PowerPoint"):
            generar_ppt(df_m, df, ventas, ganancia, margen)
            with open("reporte_final.pptx", "rb") as f:
                st.download_button("Descargar PPT", f, "reporte_final.pptx")

else:
    st.info("📂 Sube un archivo Excel")
